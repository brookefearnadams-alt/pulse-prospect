import streamlit as st
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io

# Page configuration
st.set_page_config(layout="wide", page_title="OmniPulse AI", page_icon="📡")

# Custom Styles for High Scannability
st.markdown("""
    <style>
    .main-title { font-size: 36px !important; font-weight: 800; color: #1E3A8A; margin: 0; }
    .subtitle { font-size: 15px !important; color: #4B5563; margin-bottom: 20px; font-style: italic; }
    .section-box { background-color: #F9FAFB; padding: 20px; border-radius: 8px; border: 1px solid #E5E7EB; margin-bottom: 20px; }
    .slide-red { background-color: #FEE2E2; padding: 15px; border-radius: 6px; border-left: 5px solid #EF4444; margin-bottom: 15px; }
    .slide-green { background-color: #D1FAE5; padding: 15px; border-radius: 6px; border-left: 5px solid #10B981; margin-bottom: 15px; }
    .slide-blue { background-color: #DBEAFE; padding: 15px; border-radius: 6px; border-left: 5px solid #3B82F6; margin-bottom: 15px; }
    </style>
""", unsafe_allow_html=True)

# Helper function to generate deep, multi-data PowerPoint presentation arrays (No Email, Max Intel)
def create_pptx_deck(data, clean_url):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    # Using slide layout index 6 for a clean blank canvas layout slate
    blank_layout = prs.slide_layouts[6]
    
    # ----------------------------------------------------
    # NEW INSERTION: SLIDE 0 - Institutional Title Slide
    # ----------------------------------------------------
    title_slide = prs.slides.add_slide(blank_layout)
    
    # Full Screen Deep Blue Corporate Background Card
    bg_rect = title_slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = RGBColor(30, 58, 138) # Corporate Deep Dark Blue
    bg_rect.line.fill.background()
    
    # Main Platform Title Text Container
    title_box = title_slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(1.5))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "📡 OMNIPULSE AI — ACCOUNT BRIEFING"
    p_title.font.size = Pt(40)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(255, 255, 255) # Pure White
    
    # Custom Prospect Subtitle Target Container
    sub_box = title_slide.shapes.add_textbox(Inches(1.0), Inches(3.8), Inches(11.333), Inches(2.0))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    
    p_sub1 = tf_sub.paragraphs[0]
    p_sub1.text = f"Strategic Cost-Center & Margin Diagnostic for: {data['prospect_name'].upper()}"
    p_sub1.font.size = Pt(20)
    p_sub1.font.color.rgb = RGBColor(16, 185, 129) # Emerald Green Accent Label
    p_sub1.space_after = Pt(14)
    
    p_sub2 = tf_sub.add_paragraph()
    p_sub2.text = f"Framework Allocation Mix: {data['recommended_mix']}\nPowered by Enterprise Intelligence Data & Elvex Orchestration Flows"
    p_sub2.font.size = Pt(13)
    p_sub2.font.color.rgb = RGBColor(209, 213, 219) # Light Gray Meta Copy

    # ----------------------------------------------------
    # SLIDE 1: Enterprise Data Scrape & Signal Matrix
    # ----------------------------------------------------
    slide_1 = prs.slides.add_slide(blank_layout)
    
    # Header Blue Band
    rect_1 = slide_1.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    rect_1.fill.solid()
    rect_1.fill.fore_color.rgb = RGBColor(30, 58, 138)
    rect_1.line.fill.background()
    
    txBox_1 = slide_1.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.333), Inches(0.8))
    p1 = txBox_1.text_frame.paragraphs[0]
    p1.text = f"{data['prospect_name']} — Commercial Footprint & Signal Matrix"
    p1.font.size = Pt(26)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add Large Data Grid Table Shape (Rows: 7, Columns: 2)
    table_shape = slide_1.shapes.add_table(7, 2, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    table = table_shape.table
    table.columns.width = Inches(3.5)
    table.columns.width = Inches(8.833)
    
    infra_label = "Google Floodlight Enabled Tracker | TLS 1.3 Certified Checkout" if "city" in clean_url or "furniture" in clean_url else "Klaviyo Direct Event Sync Active | TLS 1.3 Certified Flow"
    
    matrix_data = [
        ["Market Classification & Vertical", data['vertical']],
        ["Est. Monthly Digital Capital Deployment", data['estimated_digital_monthly']],
        ["Active Social/Programmatic Creative Footprint", data['meta_ad_count']],
        ["Search & Core Automation Channels", data['google_ad_types']],
        ["Verified Tracking Infrastructure Found", f"{data['pixel_detections']} | {infra_label}"],
        ["Market Share Erosion & Conquest Risks", data['competitive_threat']],
        ["Core Creative Asset Optimization Deficit", data['creative_gap']]
    ]
    
    for row_idx, row_content in enumerate(matrix_data):
        for col_idx, cell_text in enumerate(row_content):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(55, 65, 81)
            if col_idx == 0:
                p.font.bold = True
                p.font.color.rgb = RGBColor(30, 58, 138)
                
    # ----------------------------------------------------
    # SLIDE 2: Consultative Strategic Action Briefing (NO OUTREACH - MAX DATA)
    # ----------------------------------------------------
    slide_2 = prs.slides.add_slide(blank_layout)
    
    # Header Blue Band
    rect_2 = slide_2.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    rect_2.fill.solid()
    rect_2.fill.fore_color.rgb = RGBColor(30, 58, 138)
    rect_2.line.fill.background()
    
    txBox_2 = slide_2.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.333), Inches(0.8))
    p2 = txBox_2.text_frame.paragraphs[0]
    p2.text = f"{data['prospect_name']} — Consultative Optimization Index"
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255, 255, 255)
    
    # Left Column Container - Operational Strategy & Allocation
    leftBox = slide_2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.5))
    ltf = leftBox.text_frame
    ltf.word_wrap = True
    
    lp1 = ltf.paragraphs[0]
    lp1.text = "📊 Operational Strategy & Allocation"
    lp1.font.size = Pt(18)
    lp1.font.bold = True
    lp1.font.color.rgb = RGBColor(30, 58, 138)
    lp1.space_after = Pt(14)
    
    lp2_h = ltf.add_paragraph()
    lp2_h.text = "• Target Cost-Center Optimization Framework:"
    lp2_h.font.size = Pt(13)
    lp2_h.font.bold = True
    lp2_h.font.color.rgb = RGBColor(30, 58, 138)
    
    lp2_b = ltf.add_paragraph()
    lp2_b.text = data['vbr_statement']
    lp2_b.font.size = Pt(11)
    lp2_b.font.color.rgb = RGBColor(55, 65, 81)
    lp2_b.space_after = Pt(14)
    
    lp3_h = ltf.add_paragraph()
    lp3_h.text = "• Identified Capital Allocation Leak:"
    lp3_h.font.size = Pt(13)
    lp3_h.font.bold = True
    lp3_h.font.color.rgb = RGBColor(30, 58, 138)
    
    lp3_b = ltf.add_paragraph()
    lp3_b.text = data['vulnerability']
    lp3_b.font.size = Pt(11)
    lp3_b.font.color.rgb = RGBColor(55, 65, 81)
    lp3_b.space_after = Pt(14)
    
    lp4_h = ltf.add_paragraph()
    lp4_h.text = "• Proposed Marketing Capital Mix Reallocation:"
    lp4_h.font.size = Pt(13)
    lp4_h.font.bold = True
    lp4_h.font.color.rgb = RGBColor(30, 58, 138)
    
    lp4_b = ltf.add_paragraph()
    lp4_b.text = data['recommended_mix']
    lp4_b.font.size = Pt(11)
    lp4_b.font.color.rgb = RGBColor(55, 65, 81)
    
    # Right Column Container - Multi-Platform Content Pillars (Tightened Vertically)
    rightBox = slide_2.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.6))
    rtf = rightBox.text_frame
    rtf.word_wrap = True
    
    rp1 = rtf.paragraphs[0]
    rp1.text = "🎯 Executive Framework Pillars"
    rp1.font.size = Pt(18)
    rp1.font.bold = True
    rp1.font.color.rgb = RGBColor(30, 58, 138)
    rp1.space_after = Pt(8)
    
    # Pillar 1
    rp2 = rtf.add_paragraph()
    rp2.text = "1. Auction Pressures & Digital Bottlenecks"
    rp2.font.bold = True
    rp2.font.size = Pt(13)
    rp2.font.color.rgb = RGBColor(30, 58, 138)
    rp2.space_after = Pt(2)
    for b in data['slide_1_bullets']:
        p_sub = rtf.add_paragraph()
        p_sub.text = f"   - {b}"
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = RGBColor(55, 65, 81)
        p_sub.space_after = Pt(2)
    p_sub.space_after = Pt(8)
        
    # Pillar 2
    rp3 = rtf.add_paragraph()
    rp3.text = "2. Cross-Platform Portfolio Multipliers"
    rp3.font.bold = True
    rp3.font.size = Pt(13)
    rp3.font.color.rgb = RGBColor(30, 58, 138)
    rp3.space_after = Pt(2)
    for b in data['slide_2_bullets']:
        p_sub = rtf.add_paragraph()
        p_sub.text = f"   - {b}"
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = RGBColor(55, 65, 81)
        p_sub.space_after = Pt(2)
    p_sub.space_after = Pt(8)
        
    # Pillar 3
    rp4 = rtf.add_paragraph()
    rp4.text = "3. Efficiency Blueprint & Test-Zone Targets"
    rp4.font.bold = True
    rp4.font.size = Pt(13)
    rp4.font.color.rgb = RGBColor(30, 58, 138)
    rp4.space_after = Pt(2)
    for b in data['slide_3_bullets']:
        p_sub = rtf.add_paragraph()
        p_sub.text = f"   - {b}"
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = RGBColor(55, 65, 81)
        p_sub.space_after = Pt(2)
        
    binary_output = io.BytesIO()
    prs.save(binary_output)
    binary_output.seek(0)
    return binary_output

# Top Header Layout
st.markdown("<p class='main-title'>📡 OMNIPULSE AI</p>", unsafe_allow_html=True)
st.markdown("<p class='subtitle'>The Official Media Sales Intelligence Platform — Powered by Enterprise Data & Elvex Flows</p>", unsafe_allow_html=True)

# Sidebar System Navigation
st.sidebar.markdown("### 🖥️ Platform Controls")
st.sidebar.success("● Core Engine: Connected")
st.sidebar.success("### 🧠 Elvex Workflow Intelligence Nodes")
st.sidebar.info("✔ Web Scraper Node [Active]")
st.sidebar.info("✔ Cross-Platform Grounding Database [Active]")
st.sidebar.info("✔ Business Rule Formulation Persona [Active]")
st.sidebar.markdown("---")
st.sidebar.markdown("**Demo Input Mapping:**")
st.sidebar.write("• Type **'cityfurniture'** for South Florida regional retail intelligence.")
st.sidebar.write("• Type **'ecorest'** for premium DTC mattress metrics.")

# Session State Initialization
if "submitted" not in st.session_state:
    st.session_state.submitted = False

# Performance Tickers
t1, t2, t3, t4 = st.columns(4)
t1.metric(label="Average Cross-Platform CAC Drop", value="-24%")
t2.metric(label="Multi-Screen Conversion Lift", value="3.2x")
t3.metric(label="Showroom Foot Traffic Lift", value="+22%")
t4.metric(label="Elvex Processing Time", value="1.2 Sec")

st.markdown("###")

# User Input Panel
st.markdown("### 🔍 Enterprise Cost-Center & Margin Diagnostic")

if "url_input" not in st.session_state:
    st.session_state.url_input = "https://cityfurniture.com"

prospect_url = st.text_input(
    "Input Target Prospect Website URL:", 
    value=st.session_state.url_input,
    key="url_box"
)

# Button Layout alignment
btn_col1, btn_col2, _ = st.columns(3)
with btn_col1:
    if st.button("Generate Strategic Financial Briefing", type="primary"):
        st.session_state.submitted = True
with btn_col2:
    if st.button("Reset Form"):
        st.session_state.submitted = False
        st.session_state.url_input = ""
        st.rerun()

# Elvex-Enriched Strategy Database with Full Dark Funnel Suite Included
mock_database = {
    "cityfurniture": {
        "prospect_name": "CITY Furniture",
        "vertical": "B2C Retail & Home Goods / Regional Chain",
        "corporate_intel": "CEO: Andrew Koenig | HQ: Tamarac, FL | Annual Rev: \$919.5M",
        "estimated_digital_monthly": "\$250,000 - \$400,000 (Estimated Scale Spend Across South Florida Market)",
        "meta_ad_count": "High-Volume Crowded Auctions (Meta, Google PMax, Programmatic Display)",
        "google_ad_types": "Google Performance Max (PMax), Local Inventory Ads, YouTube Pre-Roll",
        "pixel_detections": "Meta Advanced Conversion API, Google Floodlight, Criteo Retargeting Tag",
        "creative_gap": "Dwell time on display assets is under 1.5 seconds. 92% of active creatives focus on margin-diluting discount promo codes instead of brand value storytelling.",
        "competitive_threat": "Rooms To Go, Ashley, Wayfair, and Amazon are running aggressive keyword conquesting ads, artificially inflating bottom-of-funnel bidding auction CPMs.",
        "vulnerability": "Paying premium acquisition rates to harvest intent that someone else built. Operating without a broad-market broadcast layer traps their budget inside volatile digital auctions.",
        "vbr_statement": "Verified cross-platform retail tracking confirms that anchoring localized digital search spend with high-frequency daytime and local news broadcast media blocks drops blended digital Customer Acquisition Cost (CAC) by 24% while driving an immediate 15% surge in high-intent organic brand searches.",

        "recommended_mix": "Convergent Mix: 50% Mass Local Linear TV Authority Layers | 30% Hyper-Targeted Connected TV | 20% Search/Social Retargeting Shield",
        "email_subject": "Optimizing CITY Furniture's customer acquisition costs and showroom foot traffic by 22%",
        "email_body": "Hi Team CITY Furniture,\n\nI’ve been tracking your market footprint across South Florida. Our enterprise intelligence analytics engine flags that your digital acquisition layer is heavily leveraged within automated ad auctions, facing intense local keyword conquesting from national scale competitors like Wayfair and Ashley.\n\nWhile your e-commerce retargeting captures immediate intent, a digital-only approach leaves you vulnerable: auction competition is inflating your acquisition costs, and average display ad dwell time has dropped below 1.5 seconds.\n\nOur cross-platform data reveals that regional retailers who anchor their search and social engines with structured daytime and local news media blocks create an insulation layer that reduces digital CAC by 24%, while triggering a 22% increase in verified showroom foot traffic.\n\nI’ve mapped out a 3-slide efficiency blueprint showing how a convergent media allocation will protect CITY Furniture's product margins this quarter. Do you have 10 minutes this Thursday at 2 PM to look over the metrics?\n\nBest,\n[Account Executive Name]",
        "slide_1_bullets": [
            "The Auction Trap: Bidding directly against multi-billion dollar national conglomerates on hyper-crowded search keywords.",
            "Margin Pressure: Digital acquisition costs (CAC) are rising as local programmatic auction bids inflate across Florida markets.",
            "The Attention Gap: Dwell time on transactional digital banner inventory sits under 1.5 seconds, failing to capture lasting household consideration."
        ],
        "slide_2_bullets": [
            "The Portfolio Multiplier: Moving away from siloed digital tactics and adding an offline mass authority layer triggers a 3.2x multi-screen conversion rate lift.",
            "The Cost Shield: Anchoring digital search and social spending with daytime and early news broadcast flights drops digital CAC by 24%.",
            "Showroom Asset Lift: Media convergence builds market-wide brand equity, yielding a proven 22% increase in verified in-store foot traffic."
        ],
        "slide_3_bullets": [
            "Convergent Reallocation: Moving budget into a resilient allocation of 50% broad-market authority blocks, 30% zip-code targeted CTV, and 20% digital retargeting capture.",
            "Advanced Attribution: Measuring campaign success through blended Marketing Efficiency Ratios (MER) and post-air organic spikes, bypassing broken click tracking pixels.",
            "Strategic Execution: Coordinating localized media test windows across high-propensity South Florida homeowner demographics next month."
        ]
    },
    "ecorest": {
        "prospect_name": "EcoRest Bedding",
        "vertical": "DTC Home Goods / Premium Mattress",
        "corporate_intel": "DTC E-Commerce Scale Player | Focus: Premium Organic Latex Lines",
        "estimated_digital_monthly": "\$45,000 - \$60,000",
        "meta_ad_count": "34 Active Creative Variations",
        "google_ad_types": "Performance Max (Heavy Search & Shopping focus)",
        "pixel_detections": "Meta Pixel, Google Tag Manager, TikTok Pixel, Klaviyo Tracking",
        "creative_gap": "Over-indexing on text-heavy testimonial static ads. Complete absence of high-production video storytelling or lifestyle video formats to build brand equity.",
        "competitive_threat": "Casper and Purple Mattresses are currently running localized geotargeted digital conquests in this market, outbidding EcoRest for the keyword 'organic latex mattress' by 42%.",
        "vulnerability": "Heavy reliance on Meta and Google Ads is causing a 38% YoY inflation in their Digital Customer Acquisition Costs (CAC), severely squeezing net margins on their organic latex mattress line.",
        "vbr_statement": "Internal campaign data proves that local DTC brands anchoring active digital campaigns with a high-impact Linear TV daytime flight experience an average 24% reduction in overall Digital CAC via a 15% surge in high-intent organic brand searches.",
        "recommended_mix": "60% Daytime Linear (Homeowner index) | 20% Connected TV (CTV) Retargeting | 20% Localized Digital News Takeovers",
        "email_subject": "Optimizing EcoRest's digital acquisition efficiency by 24%",
        "email_body": "Hi Team EcoRest,\n\nI’ve been tracking your impressive digital presence for your organic latex mattress line. Our intelligence engine flagged that you are currently managing 34 active creative variations on Meta, heavily leaning into Performance Max channels.\n\nWhile this captures baseline demand, running a digital-only acquisition strategy right now means you are highly exposed. Competitors like Casper are currently outbidding you by 42% locally on key search phrases, driving up your digital CAC by 38% YoY.\n\nOur cross-platform data shows that local brands pairing active digital with strategic daytime Linear TV see an average 24% drop in digital CAC. This builds a protective 'brand shield' around your search keywords.\n\nI’ve put together a 3-slide efficiency blueprint tailored to your active pixel profiles. Do you have 10 minutes this Thursday at 2 PM to review our data?\n\nBest,\n[Account Executive Name]",
        "slide_1_bullets": [
            "Digital Inflation: EcoRest is absorbing a 38% YoY increase in Meta/Google acquisition costs.",
            "Keyword Conquesting: Competitors are outbidding EcoRest by 42% on high-intent local search phrases.",
            "Creative Bottleneck: 34 active Meta ads rely entirely on static testimonials, ignoring screen-based video storytelling."
        ],
        "slide_2_bullets": [
            "The Halo Effect: Daytime broadcast presence drives immediate mobile/desktop organic brand searches.",
            "CAC Reduction: Proven 24% decrease in digital acquisition costs by stabilizing competitive auction dependency.",
            "Audience Capture: Reaching premium 35-64 homeowners who own 80% of local premium home assets."
        ],
        "slide_3_bullets": [
            "Efficiency Mix: Deploying a 60% Daytime Linear flight paired with high-impact CTV retargeting.",
            "Pixel Sync: Utilizing existing Meta pixels to retarget consumers exposed to our local broadcast windows.",
            "Next Steps: Reviewing customized local zone maps on Thursday."
        ]
    }
}

# --- THE DATA PROCESSING AREA ---
if st.session_state.submitted:
    with st.spinner("Processing economic and digital footprint metrics..."):
        time.sleep(1.0)
        
    clean_url = prospect_url.lower()
    
    if "city" in clean_url or "furniture" in clean_url:
        data = mock_database["cityfurniture"]
    elif "ecorest" in clean_url or "bedding" in clean_url:
        data = mock_database["ecorest"]
    else:
        data = mock_database["cityfurniture"]
        
    st.markdown("---")
    st.success(f"⚡ STRATEGIC ACCOUNT BRIEFING COMPILED FOR: {data['prospect_name'].upper()}")
    st.write(f"**Target Account:** {data['prospect_name']}  |  **Profile Intel:** {data['corporate_intel']}  |  **Analytical Framework:** Cost-Center Optimization Index")
    
    # 1. Advanced Signal Scrape Matrix Panel
    st.markdown("<div class='section-box'>", unsafe_allow_html=True)
    st.markdown(f"### 🕵️‍♂️ Advanced Signal Scrape Matrix: {data['prospect_name'].upper()}")
    
    m1, m2, m3, m4 = st.columns(4)
    with m1:
        st.markdown("**💰 Monthly Capital Allocation**")
        st.text_area("Budget", value=data['estimated_digital_monthly'], height=68, label_visibility="collapsed", disabled=True)
    with m2:
        st.markdown("**📱 Active Social Creative Sets**")
        st.text_area("Social Footprint", value=data['meta_ad_count'], height=68, label_visibility="collapsed", disabled=True)
    with m3:
        st.markdown("**🎯 Search Automation Engine**")
        st.text_area("Google Channels", value=data['google_ad_types'], height=68, label_visibility="collapsed", disabled=True)
    with m4:
        st.markdown("**⚖️ Market Share Risk Index**")
        risk_label = "HIGH VOLATILITY" if "city" in clean_url or "furniture" in clean_url else "ELEVATED RISK"
        st.text_area("Risk Index", value=risk_label, height=68, label_visibility="collapsed", disabled=True)
        
    st.markdown("---")
    
    col_sig1, col_sig2 = st.columns(2)
    with col_sig1:
        st.markdown("#### 💻 Technical Infrastructure Diagnostics")
        st.markdown(f"• **Active Tracking Pixels Verified:** `{data['pixel_detections']}`")
        st.markdown("• **Attribution Tracking Script:** `Google Floodlight Tracker Active`" if "city" in clean_url or "furniture" in clean_url else "• **Attribution Tracking Script:** `Klaviyo Event Tracking Active`")
        st.markdown("• **SSL Security Framework:** `TLS 1.3 Certified (Secure Checkout Flow)`")
        st.markdown("• **Cross-Channel Integration Anchor:** `Elvex Flow API Payload Sync Complete`")
        
    with col_sig2:
        st.markdown("#### 📉 Strategic Gaps & Market Pressures")
        st.markdown(f"⚠️ **Market Share Erosion Threat:** {data['competitive_threat']}")
        st.markdown(f"💡 **Creative Asset Gap:** {data['creative_gap']}")
        st.markdown("• **Estimated Digital Share-of-Voice (SOV) Slip:** `Drop 14% YoY due to National Bidding Pods`" if "city" in clean_url or "furniture" in clean_url else "• **Estimated Digital Share-of-Voice (SOV) Slip:** `Drop 22% YoY vs Venture-Backed Competitors`")
        
    st.markdown("</div>", unsafe_allow_html=True)
    
        # NEW ELEMENT: Elvex Dark Funnel Visibility Matrix (Fully Expanded Research Suite)
    st.markdown("<div class='section-box' style='border-left: 5px solid #7C3AED;'>", unsafe_allow_html=True)
    st.markdown("### 🔮 Elvex Dark Funnel Visibility Matrix")
    
    # Macro Status Indicators
    df_col1, df_col2 = st.columns(2)
    with df_col1:
        st.markdown("#### 🚫 Current Digital Attribution Blind-Spots")
        st.markdown("• **Dark Social Leakage:** `Estimated 68% of sectional research happens on untracked messaging apps (Slack/WhatsApp).`" if "city" in clean_url or "furniture" in clean_url else "• **Dark Social Leakage:** `Estimated 74% of mattress review link sharing happens via dark messaging nodes.`")
        st.markdown("• **Last-Click Fallacy Index:** `High. Analytics tools misattribute offline brand-equity loops directly to organic search or direct web traffic keywords.`")
        st.markdown("• **Unmeasured Touchpoints:** `Podcasts, community recommendations, and video impressions are completely hidden from current dashboard view graphs.`")
        
    with df_col2:
        st.markdown("#### 💡 The Broadcast Attribution Fix")
        st.markdown("• **Blended MER Tracking Optimization:** `Capturing aggregate revenue lift patterns across the total advertising pipeline.`")
        st.markdown("• **Post-Air Organic Spike Correlation:** `Mapping baseline web navigation surges within an 8-minute window of broadcast ad exposures.`")
        st.markdown("• **Unified Brand Insulation Index:** `Using broadcast weight to build stable consumer demand outside of volatile digital auctions.`")
        
    st.markdown("---")
    
    # Deep-Dive Dark Funnel Research Sub-Panels
    st.markdown("#### 🛋️ Deep Audit Discovery: CITY Furniture Segment Data" if "city" in clean_url or "furniture" in clean_url else "#### 🛏️ Deep Audit Discovery: EcoRest Bedding Segment Data")
    
    sub_col1, sub_col2 = st.columns(2)
    with sub_col1:
        st.markdown("##### 📈 Strategic Market Strengths")
        st.write("• **Brand Visibility (Local Familiarity Index):** High baseline recognition due to physical showroom footprint and ongoing promo cycles. Crucial because furniture purchases are infrequent and consumers default to names they recognize.")
        st.write("• **Search Presence Authority:** Strong branded search volumes ('City Furniture near me') alongside high-intent category phrases (sectional, sofa, bedroom set).")
        st.write("• **Trust Signals Baseline:** Benefits from perceived legitimacy metrics: delivery capability, physical footprint, financing ease, and warranties.")
        
        st.markdown("##### 🛒 4-Stage Strategic Customer Journey")
        st.write("**1. Awareness:** Triggered by life events (moving, back pain, redecorating). Gaps exist in educational/inspiration video content outside of price promos.")
        st.write("**2. Research:** Googling non-branded help-me-choose terms. Local pack reputation is uneven across regional zip codes.")
        st.write("**3. Consideration:** Comparing 2–5 regional showrooms. Retargeting loops fail to handle core tracking objections like pet-proofing, setup care, or comfort metrics.")
        st.write("**4. Decision:** Driven by final last-mile confidence. Reassurance indicators regarding tracking windows or damage resolutions are missing at checkout.")
        
    with sub_col2:
        st.markdown("##### 📉 Identified Operational Leaks & Gaps")
        st.write("• **Mid-Funnel Content Deficit:** Web assets over-index on product grids but heavily leak traffic on research queries like sectional sizing guides, performance fabric comparisons, and room-planning software.")
        st.write("• **Review Asset Under-Utilization:** Location reviews exist across individual local showrooms but are completely un-packaged into digital ad creative hooks, testimonial reels, or dedicated landing pages.")
        st.write("• **Word-of-Mouth Capture Failure:** Zero systematic triggers or post-delivery referral requests deployed inside the critical 24–72 hour customer satisfaction window.")
        
        st.markdown("##### 🏆 Recommended Strategic Interventions")
        st.write("• **YouTube/CTV 'Familiarity' Flight:** Re-establishing upper-funnel baseline equity outside platform auction loops to lower search costs by 24%.")
        st.write("• **Objection-Handling Retargeting Module:** Shifting programmatic assets away from raw discounts to resolve uncertainties about last-mile setups, financing, or durability.")
        st.write("• **Local SEO Reputation Optimization:** Executing systematic store-level reputation updates to clean up Map Pack trust scores across regional territories.")
        
    st.markdown("</div>", unsafe_allow_html=True)

    # 2. Strategy & Pitch Text Blocks (Upgraded to multi-tabbed research suite)
    st.markdown("### 🎯 Consultative Account Strategy & Deep Audit Suite")
    
    # Draw native interactive tab headers across your center UI pane
    tab_core, tab_journey, tab_trust, tab_tracks = st.tabs([
        "📊 Core Strategy Brief", 
        "🔮 Customer Journey Audit", 
        "🛠️ Trust-Building Plan", 
        "🎙️ Seller Talk Tracks"
    ])
    
    with tab_core:
        st.info(f"**The Valid Business Reason (VBR):**\n\n{data['vbr_statement']}")
        st.error(f"**Identified Operational Inefficiency:**\n\n{data['vulnerability']}")
        if "dark_funnel_summary" in data:
            st.warning(f"**Dark Funnel Strategic Summary:**\n\n{data['dark_funnel_summary']}")
            
    with tab_journey:
        st.markdown("#### 🛒 4-Stage Customer Journey Breakdown & Tracking Gaps")
        st.caption("Surfacing hidden pre-click discovery leaks where standard conversion tracking pixels are blind.")
        
        st.markdown(f"**1. Awareness Phase (Trigger Realization):**\n*{data.get('journey_awareness', 'N/A')}*")
        st.markdown("---")
        st.markdown(f"**2. Research Phase (Non-Branded Discovery):**\n*{data.get('journey_research', 'N/A')}*")
        st.markdown("---")
        st.markdown(f"**3. Consideration Phase (Objection Comparison):**\n*{data.get('journey_consideration', 'N/A')}*")
        st.markdown("---")
        st.markdown(f"**4. Decision Phase (Last-Mile Reassurance):**\n*{data.get('journey_decision', 'N/A')}*")
        
    with tab_trust:
        st.markdown("#### 🛡️ Omni-Channel Brand Trust Optimization Plan")
        if "trust_plan" in data:
            st.success(data['trust_plan'])
        else:
            st.info("General trust framework applied.")
            
    with tab_tracks:
        st.markdown("#### 📢 Consultative C-Suite Seller Talk Tracks")
        if "seller_talk_track" in data:
            st.write(data['seller_talk_track'])
            st.caption("💡 * AE Strategy Pitch Tip: Deliver this tracking context directly to the CFO/CMO to position broadcast media as an efficiency multiplier for their active digital spend.*")
        else:
            st.write("Standard consultative positioning tracks loaded.")

    
    # 3. Outreach Copy Block
    st.markdown("### 📨 Automated Executive Outreach Script")
    st.text_input("Recommended Subject Line:", data['email_subject'])
    st.text_area("Generated Consultative Copy (Ready to Copy/Paste):", data['email_body'], height=250)
    
    # 4. Presentation Slides Content Header with Real PPTX Downloader
    st.markdown("---")
    
    dl_col1, dl_col2 = st.columns(2)
    with dl_col1:
        st.markdown("### 📊 Executive Presentation Content Architecture")
        st.write(f"**Proposed Marketing Capital Reallocation Mix:** {data['recommended_mix']}")
    with dl_col2:
        # Pass clean_url into the upgraded PowerPoint layout engine
        pptx_buffer = create_pptx_deck(data, clean_url)
        st.download_button(
            label="📊 Download PowerPoint Slides (.PPTX)",
            data=pptx_buffer.getvalue(),
            file_name=f"{data['prospect_name'].lower().replace(' ', '_')}_executive_deck.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            type="secondary"
        )

        
    st.caption("Utilize these structured business cases directly inside your presentation deck to execute a direct consultative sale:")
    
    # Slide 1 Visual
    st.markdown("<div class='slide-red'>", unsafe_allow_html=True)
    st.markdown(f"**SLIDE 1: {data['prospect_name']} — The Digital Bottleneck (Auction Scaling Pressures)**")
    for bullet in data['slide_1_bullets']:
        st.markdown(f"❌ {bullet}")
    st.markdown("</div>", unsafe_allow_html=True)
        
    # Slide 2 Visual
    st.markdown("<div class='slide-green'>", unsafe_allow_html=True)
    st.markdown("**SLIDE 2: The Cross-Platform Portfolio Multiplier Effect (3.2x Conversion Lift)**")
    for bullet in data['slide_2_bullets']:
        st.markdown(f"📈 {bullet}")
    st.markdown("</div>", unsafe_allow_html=True)
        
    # Slide 3 Visual
    st.markdown("<div class='slide-blue'>", unsafe_allow_html=True)
    st.markdown("**SLIDE 3: The Efficiency Blueprint (Capital Reallocation and Margin Protection)**")
    for bullet in data['slide_3_bullets']:
        st.markdown(f"🛠️ {bullet}")
    st.markdown("</div>", unsafe_allow_html=True)
